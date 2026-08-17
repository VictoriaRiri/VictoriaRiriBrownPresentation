#!/usr/bin/env python3
"""
make_presentation.py
Generates Victoria_Riri_Wagura_Brown_Application.pptx using portrait.jpg present in the repo root.
Requires: python-pptx, pillow, matplotlib
This script is non-interactive and designed to run in CI (GitHub Actions) or locally.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

OUT_PPTX = "Victoria_Riri_Wagura_Brown_Application.pptx"
PORTRAIT = "portrait.jpg"

# Fallback: pick any image in the repo if portrait.jpg not present
if not os.path.exists(PORTRAIT):
    imgs = [f for f in os.listdir('.') if f.lower().endswith(('.jpg', '.jpeg', '.png'))]
    if imgs:
        PORTRAIT = imgs[0]
    else:
        raise SystemExit("No image file found in repo root. Please add portrait.jpg and re-run.")

# Stat data for chart
years = ["2021", "2022", "2023", "2024"]
counts = [3, 3, 2, 2]

# Colors (Brown palette)
BROWN = RGBColor(78, 42, 30)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)


def set_paragraph_format(tx, font_name="Open Sans", size=18, bold=False, color=BROWN):
    for paragraph in tx.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_text = notes_slide.notes_text_frame
    notes_text.clear()
    notes_text.text = text

# 1 — Cover
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
left = top = Inches(0)
width = prs.slide_width
height = prs.slide_height
# Background: use portrait cropped to slide aspect ratio
img = Image.open(PORTRAIT)
img = img.convert("RGB")
img_w, img_h = img.size
slide_aspect = prs.slide_width / prs.slide_height
img_aspect = img_w / img_h
if img_aspect > slide_aspect:
    new_w = int(slide_aspect * img_h)
    left_crop = (img_w - new_w)//2
    img = img.crop((left_crop, 0, left_crop + new_w, img_h))
else:
    new_h = int(img_w / slide_aspect)
    top_crop = (img_h - new_h)//2
    img = img.crop((0, top_crop, img_w, top_crop + new_h))
img.save("cover_bg.jpg", quality=85)
slide.shapes.add_picture("cover_bg.jpg", left, top, width=width, height=height)
# Title
tx_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12.1), Inches(1.4))
tf = tx_box.text_frame
tf.text = "Victoria Riri Wagura — Why Brown"
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
p.runs[0].font.size = Pt(40)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = RGBColor(255,255,255)
# Subtitle
sub = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.6))
st = sub.text_frame
st.text = "Applicant | Module 2"
st.paragraphs[0].runs[0].font.size = Pt(20)
st.paragraphs[0].runs[0].font.color.rgb = RGBColor(255,255,255)
add_notes(slide, "Quick greet: Hello — I’m Victoria Riri Wagura. This short deck explains why Brown’s Open Curriculum and research environment are the natural next step for my interdisciplinary work in perception and systems engineering.")

# 2 — About slide (photo + bullets)
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.shapes.add_picture(PORTRAIT, Inches(0.6), Inches(0.6), width=Inches(4.0), height=Inches(5.3))
tx = slide.shapes.add_textbox(Inches(5.2), Inches(0.8), Inches(7.0), Inches(5.0)).text_frame
tx.text = "About"
tx.paragraphs[0].runs[0].font.size = Pt(32)
p = tx.add_paragraph(); p.text = "Name: Victoria Riri Wagura"; p.level = 1
p = tx.add_paragraph(); p.text = "Home country: Kenya"; p.level = 1
p = tx.add_paragraph(); p.text = "Module: 2"; p.level = 1
p = tx.add_paragraph(); p.text = "Intended major: Computer Engineering / Robotics (flexible)"; p.level = 1
set_paragraph_format(tx, font_name="Open Sans", size=18, bold=False, color=BROWN)
add_notes(slide, "Short personal context: I build systems that connect people and machines. My background spans deployed products and perception research prototypes — I care about engineering that serves people.")

# 3 — Stats chart slide
plt.figure(figsize=(6.4,3.6))
bars = plt.bar(years, counts, color=["#4E2A1E"])
plt.title("Enrollment from Kenya — recent years")
plt.ylabel("Number of students")
plt.ylim(0, max(counts)+1)
for i, v in enumerate(counts):
    plt.text(i, v + 0.05, str(v), ha='center', fontsize=12)
chartfile = "enrollment_chart.png"
plt.tight_layout()
plt.savefig(chartfile, dpi=150)
plt.close()
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.shapes.add_picture(chartfile, Inches(1.2), Inches(1.0), width=Inches(10))
tx = slide.shapes.add_textbox(Inches(1.2), Inches(5.0), Inches(10), Inches(1)).text_frame
tx.text = "Counts reflect enrolled students from Kenya. Source: Francis Parserio (LinkedIn)."
add_notes(slide, "This small but stable set of students from Kenya highlights how each admits' experience matters. I’m citing Francis Parserio’s note on enrollment trends (link in Sources).")

# 4 — Why Brown (100 words)
why_text = ("Brown’s Open Curriculum is the environment I need to grow as a generalist who pursues deep technical work across multiple domains. "
            "I’m one of those learners with many loves: a curiosity for movement perception, a builder’s instinct for systems integration, "
            "and a commitment to technology that serves communities. At Brown I can explore these threads freely — testing ideas through Satisfactory/No Credit "
            "while discovering whether a subject is a passing interest or a true intellectual home. That freedom, paired with Brown’s culture of collaborative scholarship, "
            "is the reason I want four focused years to turn exploratory projects into rigorous research and impact.")
slide = prs.slides.add_slide(prs.slide_layouts[6])
tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.6), Inches(5.5)).text_frame
tx.text = "Why Brown"
tx.paragraphs[0].runs[0].font.size = Pt(30)
p = tx.add_paragraph()
p.text = why_text
p.level = 1
p.runs[0].font.size = Pt(16)
set_paragraph_format(tx, font_name="Open Sans", size=16, bold=False, color=BROWN)
add_notes(slide, "Expand briefly on personal examples: MotionTrace (perception), Jewellery818 (systems), Smartsolve (human-centered service).")

# 5 — Academic interests
slide = prs.slides.add_slide(prs.slide_layouts[6])
tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.6), Inches(5.5)).text_frame
tx.text = "Academic interests"
tx.paragraphs[0].runs[0].font.size = Pt(28)
p = tx.add_paragraph()
p.text = "• Intended major: Computer Engineering / Robotics (open to interdisciplinary pathway)"
p.level = 1
p = tx.add_paragraph(); p.text = "• Concentrations: perception & human-robot interaction, systems engineering, AI & society"; p.level = 1
p = tx.add_paragraph(); p.text = "• Course types: advanced robotics, probabilistic perception, ethics & society seminars"; p.level = 1
set_paragraph_format(tx, font_name="Open Sans", size=18, bold=False, color=BROWN)
add_notes(slide, "Mention eagerness to take advantage of cross‑department advising; intent to work with labs and take S/NC to explore.")

# 6 — Projects
slide = prs.slides.add_slide(prs.slide_layouts[6])
tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12.0), Inches(6.5)).text_frame
tx.text = "Project portfolio — how the pieces connect"
tx.paragraphs[0].runs[0].font.size = Pt(26)
p = tx.add_paragraph(); p.text = "MotionTrace — pose estimation and perception. MotionTrace extracts structure from human movement; it’s the perception problem that lets machines see bodies."; p.level = 1
p = tx.add_paragraph(); p.text = "Jewellery818 — full-stack systems engineering & deployment. Integration: CMS, JSONBin sync, Vercel; built the muscle of product engineering."; p.level = 1
p = tx.add_paragraph(); p.text = "Smartsolve — free tool for Kenyan students; demonstrates human-centered motive and community focus."; p.level = 1
set_paragraph_format(tx, font_name="Open Sans", size=16, bold=False, color=BROWN)
add_notes(slide, "Narrative: Jewellery818 taught systems and deployment; limited traction led to MotionTrace — perception; Smartsolve shows ethics & access. This arc maps to Tellex’s lab interests.")

# 7 — Research fit: Tellex
slide = prs.slides.add_slide(prs.slide_layouts[6])
tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.6), Inches(5.5)).text_frame
tx.text = "Why Professor Tellex’s lab"
tx.paragraphs[0].runs[0].font.size = Pt(28)
p = tx.add_paragraph(); p.text = "Shared core question: before a system can respond to people, it must first perceive them accurately."; p.level = 1
p = tx.add_paragraph(); p.text = "My interest: developing algorithms to connect perception (pose, gesture) with interactive response."; p.level = 1
p = tx.add_paragraph(); p.text = "MotionTrace made me confront the same central problem Professor Tellex addresses: a system can’t respond meaningfully to a person until it can see what the person is doing."; p.level = 1
set_paragraph_format(tx, font_name="Open Sans", size=16, bold=False, color=BROWN)
add_notes(slide, "Tie to lab methods or proposed work: e.g., integrating pose outputs into dialogue and robot action pipelines.")

# 8 — Campus life & landmarks
slide = prs.slides.add_slide(prs.slide_layouts[6])
tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12.0), Inches(5.0)).text_frame
tx.text = "Campus life and landmarks"
tx.paragraphs[0].runs[0].font.size = Pt(28)
p = tx.add_paragraph(); p.text = "• Granoff Center for the Creative Arts — interdisciplinary creative space."; p.level = 1
p = tx.add_paragraph(); p.text = "• Van Wickle Gates — ceremonial campus gateway; a community gathering place."; p.level = 1
p = tx.add_paragraph(); p.text = "• Sciences Library — central STEM study resource."; p.level = 1
p = tx.add_paragraph(); p.text = "• Main Green — historical gathering place; where students heard news of the Open Curriculum."; p.level = 1
set_paragraph_format(tx, font_name="Open Sans", size=16, bold=False, color=BROWN)
add_notes(slide, "Include these places because they speak to Brown’s culture — public, collaborative, and creative.")

# 9 — News & context
slide = prs.slides.add_slide(prs.slide_layouts[6])
tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.6), Inches(5.0)).text_frame
tx.text = "News & context"
tx.paragraphs[0].runs[0].font.size = Pt(26)
p = tx.add_paragraph(); p.text = "Recent reporting by Francis Parserio highlights enrollment counts and campus context — included here for transparency and context."; p.level = 1
p = tx.add_paragraph(); p.text = "Article link (see Sources): https://share.google/yhBkQuBqDMW9zZ94v"; p.level = 1
set_paragraph_format(tx, font_name="Open Sans", size=14, bold=False, color=BROWN)
add_notes(slide, "This article is included as context for the enrollment numbers and as a public source. It does not change my application narrative; it situates the statistics.")

# 10 — Sources & credits
slide = prs.slides.add_slide(prs.slide_layouts[6])
tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.6), Inches(5.5)).text_frame
tx.text = "Sources & image credits"
tx.paragraphs[0].runs[0].font.size = Pt(26)
p = tx.add_paragraph(); p.text = "• Enrollment statistics: Francis Parserio — LinkedIn / article"; p.level = 1
p = tx.add_paragraph(); p.text = "• News article: https://share.google/yhBkQuBqDMW9zZ94v"; p.level = 1
p = tx.add_paragraph(); p.text = "• Portrait: supplied by applicant (used with permission)"; p.level = 1
p = tx.add_paragraph(); p.text = "• Campus images: use Brown University official pages or campus photos with attribution"; p.level = 1
set_paragraph_format(tx, font_name="Open Sans", size=14, bold=False, color=BROWN)
add_notes(slide, "If the application requires the official logo or color usage rules, confirm and I’ll adjust.")

# 11 — Thank you / contact
slide = prs.slides.add_slide(prs.slide_layouts[6])
tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.6), Inches(4.5)).text_frame
tx.text = "Thank you"
tx.paragraphs[0].runs[0].font.size = Pt(32)
p = tx.add_paragraph(); p.text = "Victoria Riri Wagura"; p.level = 1
p = tx.add_paragraph(); p.text = "Email: victoriawagura@gmail.com"; p.level = 1
p = tx.add_paragraph(); p.text = "I welcome the opportunity to discuss research fit and coursework."; p.level = 1
if os.path.exists(PORTRAIT):
    slide.shapes.add_picture(PORTRAIT, Inches(10.5), Inches(4.4), width=Inches(2.0), height=Inches(2.6))

prs.save(OUT_PPTX)
print("Created", OUT_PPTX)
